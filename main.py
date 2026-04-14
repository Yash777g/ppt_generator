from fastapi import FastAPI, HTTPException
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from groq import Groq

import anthropic
import json
import os
import uuid
import re

app = FastAPI()

# Mount static files (frontend)
app.mount("/static", StaticFiles(directory="static"), name="static")

# When using anthropic api key (i.e having money in hand)
# client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))

# using groq api key as it free 
client = Groq(api_key="API_key")

THEMES = {
    "professional": {
        "bg": RGBColor(0x1B, 0x2A, 0x4A),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "slide_bg": RGBColor(0xFF, 0xFF, 0xFF),
        "slide_title": RGBColor(0x1B, 0x2A, 0x4A),
        "slide_body": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0x4A, 0x90, 0xD9),
    },
    "minimal": {
        "bg": RGBColor(0xF5, 0xF5, 0xF3),
        "title_text": RGBColor(0x1A, 0x1A, 0x1A),
        "slide_bg": RGBColor(0xFF, 0xFF, 0xFF),
        "slide_title": RGBColor(0x22, 0x22, 0x22),
        "slide_body": RGBColor(0x55, 0x55, 0x55),
        "accent": RGBColor(0x88, 0x87, 0x80),
    },
    "vibrant": {
        "bg": RGBColor(0x0F, 0x6E, 0x56),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "slide_bg": RGBColor(0xFF, 0xFF, 0xFF),
        "slide_title": RGBColor(0x0F, 0x6E, 0x56),
        "slide_body": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0xD8, 0x5A, 0x30),
    },
}


class GenerateRequest(BaseModel):
    topic: str
    description: str = ""
    slide_count: int = 7
    theme: str = "professional"


def generate_slides_with_claude(topic: str, description: str, slide_count: int) -> list:
    prompt = f"""You are a professional presentation designer. Generate a complete slide deck.

Topic: {topic}
{f"Description: {description}" if description else ""}
Number of slides: {slide_count} (including title slide and conclusion)

Return ONLY a valid JSON array. No explanation, no markdown fences. Each element:
{{
  "slide": <number>,
  "title": "<slide title>",
  "bullets": ["<point 1>", "<point 2>", "<point 3>"]
}}

Rules:
- Slide 1 is the title slide (bullets = 1-2 subtitle lines)
- Last slide is "Key Takeaways" or "Conclusion"
- 3-5 bullet points per content slide
- Keep bullets concise (max 12 words each)
- Ensure logical content flow throughout"""

    message = client.chat.completions.create(
        model="llama-3.3-70b-versatile",  # upgrade or change to any model as per need
        max_tokens=2000,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = message.choices[0].message.content.strip()
    raw = re.sub(r"```json|```", "", raw).strip()
    return json.loads(raw)


def build_pptx(slides_data: list, theme_name: str, topic: str) -> str:
    theme = THEMES.get(theme_name, THEMES["professional"])
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for i, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        if i == 0:
            fill.fore_color.rgb = theme["bg"]
        else:
            fill.fore_color.rgb = theme["slide_bg"]

        # Accent bar (left side)
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0),
            Inches(0.15), Inches(7.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = theme["accent"]
        bar.line.fill.background()

        if i == 0:
            # Title slide
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_info["title"]
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = theme["title_text"]

            if slide_info.get("bullets"):
                sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1.5))
                stf = sub_box.text_frame
                stf.word_wrap = True
                sp = stf.paragraphs[0]
                sp.text = " | ".join(slide_info["bullets"])
                sp.font.size = Pt(20)
                sp.font.color.rgb = RGBColor(0xC8, 0xD8, 0xF0) if theme_name == "professional" else theme["title_text"]
        else:
            # Content slide - title
            title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(1.1))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_info["title"]
            p.font.size = Pt(30)
            p.font.bold = True
            p.font.color.rgb = theme["slide_title"]

            # Divider line
            line = slide.shapes.add_shape(1, Inches(0.4), Inches(1.4), Inches(12.5), Pt(2))
            line.fill.solid()
            line.fill.fore_color.rgb = theme["accent"]
            line.line.fill.background()

            # Bullet points
            content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.2), Inches(5.4))
            ctf = content_box.text_frame
            ctf.word_wrap = True

            for j, bullet in enumerate(slide_info.get("bullets", [])):
                if j == 0:
                    cp = ctf.paragraphs[0]
                else:
                    cp = ctf.add_paragraph()
                cp.text = f"• {bullet}"
                cp.font.size = Pt(20)
                cp.font.color.rgb = theme["slide_body"]
                cp.space_before = Pt(8)

        # Slide number (skip title)
        if i > 0:
            num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3))
            ntf = num_box.text_frame
            np_ = ntf.paragraphs[0]
            np_.text = str(i + 1)
            np_.font.size = Pt(10)
            np_.font.color.rgb = theme["accent"]
            np_.alignment = PP_ALIGN.RIGHT

    filename = f"{uuid.uuid4().hex}.pptx"
    out_path = "test.pptx"
    prs.save(out_path)
    return out_path, filename


@app.get("/")
def index():
    return FileResponse("static/index.html")


@app.post("/generate")
def generate(req: GenerateRequest):
    if not req.topic.strip():
        raise HTTPException(status_code=400, detail="Topic is required")
    if req.slide_count < 5 or req.slide_count > 15:
        raise HTTPException(status_code=400, detail="Slide count must be between 5 and 15")

    slides_data = generate_slides_with_claude(req.topic, req.description, req.slide_count)
    out_path, filename = build_pptx(slides_data, req.theme, req.topic)

    return JSONResponse({
        "filename": filename,
        "slides": slides_data,
        "download_url": f"/download/{filename}"
    })


@app.get("/download/{filename}")
def download(filename: str):
    path = f"/tmp/{filename}"
    if not os.path.exists(path):
        raise HTTPException(status_code=404, detail="File not found")
    topic_name = "presentation"
    return FileResponse(
        path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{topic_name}.pptx"
    )
